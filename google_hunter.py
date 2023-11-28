import io
import requests
import re
import sys
import argparse
from serpapi import GoogleSearch
import csv
from concurrent.futures import ThreadPoolExecutor
import asyncio
import re
import openpyxl
from io import BytesIO
from PyPDF2 import PdfReader
import docx
import pptx
import aiohttp

# Array of google dorks
google_dorks = [
    'inurl:"@example.com"',
    'cache:example.comsite:linkedin.com "@example.com"',
    'intext:"@example.com" site:facebook.com',
    'site:stackoverflow.com "@example.com"',
    'intext:@example.com email',
    'inurl:staff "@example.com"',
    '"@example.com" intext:"contact"',
    'intext:"@example.com" filetype:txt',
    'intext:"@example.com" site:tumblr.com',
    'site:example.com',
    '@example.com',
    'link:example.com',
    'inurl:"contact" "@example.com"',
    'intitle:"example.com email"',
    'intitle:"@example.com"',
    'filetype:txt "@example.com"',
    'site:medium.com "@example.com"',
    'intext:"@example.com" site:speakerdeck.com',
    'site:*.example.com inurl:contributor',
    'inurl:team "@example.com"',
    'site:*.example.com inurl:directory',
    'inurl:about "@example.com"',
    'intext:"@example.com" filetype:doc',
    'intext:"@example.com" site:vimeo.com',
    'inurl:people "@example.com"',
    'inurl:profiles "@example.com"',
    'site:*.example.com inurl:profile',
    'related:example.com',
    'site:twitter.com "@example.com" -"tweets"',
    'intext:"@example.com" site:youtube.com',
    'site:*.example.com filetype:pdf',
    'inurl:"directory" "@example.com"',
    'inurl:users "@example.com"',
    'filetype:doc "@example.com"',
    'intext:"@example.com" filetype:ppt',
    'filetype:pdf "@example.com"',
    'site:linkedin.com "@example.com"',
    'site:*.example.com inurl:emails',
    'filetype:xls "@example.com"',
    'site:*.example.com inurl:people',
    '"@example.com" intext:"email" -site:example.com',
    'intext:"@example.com" filetype:pdf',
    'inurl:alumni "@example.com"',
    'site:github.com "@example.com"',
    'inurl:"staff" "@example.com"',
    'intext:"@example.com" site:instagram.com',
    'site:*.example.com filetype:doc',
    'site:*.example.com inurl:email',
    'site:*.example.com inurl:team',
    'inurl:"about us" "@example.com"',
    'site:*.example.com inurl:member',
    'site:zoominfo.com "@example.com"',
    'site:*.example.com inurl:about',
    'site:facebook.com "@example.com"',
    'inurl:contacts "@example.com"',
    'intext:"@example.com" site:pinterest.com',
    'intext:"@example.com" site:researchgate.net',
    'inurl:directory "@example.com"',
    'intext:"@example.com" site:academia.edu',
    'site:twitter.com "@example.com"',
    'inurl:members "@example.com"',
    'intext:"@example.com" site:slideshare.net',
    'site:*.example.com inurl:contacts',
    'inurl:contact "@example.com"',
    'site:*.example.com inurl:contact',
    'intext:"@example.com" filetype:xls',
    'site:rocketreach.co "@example.com"',
    '"@example.com" -site:example.com -site:mail.example.com',
    'site:gitlab.com "@example.com"',
    '"@example.com"',
    'site:crunchbase.com "@example.com"',
    'site:*.example.com filetype:xls',
    'site:*.example.com inurl:faculty',
    'site:*.example.com filetype:txt',
    'inurl:emails "@example.com"',
    'inurl:contributors "@example.com"',
    'site:*.example.com inurl:alumni',
    'inurl:"@example.com" filetype:csv',
    'inurl:github.com "@example.com"',
    'filetype:xlsx "@example.com"',
    'filetype:pptx "@example.com"',
    'inurl:faculty "@example.com"',
    'site:*.example.com inurl:staff',
    'site:*.example.com inurl:user',
    '"@example.com" -site:example.com',
    'intext:"@example.com" site:reddit.com',
    'site:reddit.com "@example.com"',
    'site:quora.com "@example.com"',
    'allintext: @example.com'
]

# Set up command-line argument parsing
parser = argparse.ArgumentParser(description="Email Hunter: A tool to scrape emails from a given domain.")
parser.add_argument("-k", "--api_key", required=True, help="API key for SerpAPI")
parser.add_argument("-d", "--domain", required=True, help="Domain for email addresses") 
parser.add_argument("-m", "--max_results", default=100, type=int, help="Max search results")
parser.add_argument("-r", "--regex", required=True, help="Regex pattern for matching emails")
args = parser.parse_args()

# Initialize settings from arguments
DOMAIN = args.domain  
MAX_RESULTS = args.max_results
API_KEY = args.api_key
EMAIL_REGEX = args.regex

# Compiled regex patterns for efficiency

DOMAIN_REGEXs = [
    r'\b[A-Za-z0-9._%+-]+@example.com\b',  # Standard email pattern
    r'[A-Za-z0-9._%+-]+ \[at\] example.com',  # Obfuscated pattern
    r'[A-Za-z0-9._%+-]+ at example.com',  # Another obfuscated pattern
    r'[A-Za-z0-9._%+-]+[ ]?@[ ]?example.com',  # Spaces around @
    r'[A-Za-z0-9._%+-]+\(at\)example.com',  # (at) instead of @
]
EXACT_REGEX = re.compile(rf"{EMAIL_REGEX}@{re.escape(DOMAIN)}")

# Define the search terms for the Google search
search_terms = [dork.replace('example.com', DOMAIN) for dork in google_dorks]

# Initialize dictionaries to store found emails
exact_emails = {}
found_emails = {}

async def fetch_emails(session, link):
    """ Fetch emails from a link. Works with standard HTML, PDF, Word, Excel, and PowerPoint files."""
    print(f"Fetching emails from {link}...")
    try:
        async with session.get(link) as response:
            content_type = response.headers.get('Content-Type', '')

            if 'text/html' in content_type:
                print("Found an HTML file. Extracting text...")
                text = await response.text()
            elif 'application/pdf' in content_type:
                print("Found a PDF file. Extracting text...")
                bytes_io = io.BytesIO(await response.read())
                reader = PdfReader(bytes_io)
                text = ''
                for page in reader.pages:
                    text += page.extract_text()
            elif 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' in content_type:
                print("Found a Word file. Extracting text...")
                doc = docx.Document(BytesIO(await response.read()))
                text = '\n'.join([para.text for para in doc.paragraphs])
            elif 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' in content_type:
                print("Found an Excel file. Extracting text...")
                wb = openpyxl.load_workbook(BytesIO(await response.read()), data_only=True)
                text = ''
                for sheet in wb:
                    for row in sheet.iter_rows(values_only=True):
                        text += ' '.join([str(cell) for cell in row if cell is not None]) + '\n'
            elif 'application/vnd.openxmlformats-officedocument.presentationml.presentation' in content_type:
                print("Found a PowerPoint file. Extracting text...")
                ppt = pptx.Presentation(BytesIO(await response.read()))
                text = ''
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + '\n'
            else:
                print(f"Unknown file type: {content_type}")
                text = await response.text()

            print("Searching for emails...")
            all_domain_emails = []
            for domain_regex in DOMAIN_REGEXs:
                domain_regex = domain_regex.replace('example.com', DOMAIN)
                domain_emails = re.findall(domain_regex, text)
                all_domain_emails.extend(domain_emails)

            exact_match_emails = [email for email in all_domain_emails if re.match(EXACT_REGEX, email)]

            exact = organize_emails(exact_match_emails, link)
            found = organize_emails(all_domain_emails, link)
            print(f"Found {len(exact_match_emails)} exact matches and {len(all_domain_emails)} emails for {link}")
            return exact, found
    except Exception as e:
        print(f"Error fetching {link}: {e}")
        return {}, {}

def organize_emails(emails, domain):
    """ Organize emails into a dictionary. """
    organized_emails = {}
    for email in emails:
        if email not in organized_emails:
            organized_emails[email] = set()
        organized_emails[email].add(domain)
    return organized_emails

def write_emails_to_csv(filename, email_dict):
    """ Write emails to a CSV file. """
    with open(filename, 'w', newline='') as f:
        writer = csv.writer(f)
        writer.writerow(['Email', 'Domains'])

        for email, domains in email_dict.items():
            writer.writerow([email, ', '.join(domains)])


async def main_async():
    print(f"Starting email search for domain {DOMAIN}")
    print(f"Using regex pattern: {EMAIL_REGEX}")
    print(f"Maximum Google search results per query: {MAX_RESULTS}")
    print("Beginning search queries...")


    async with aiohttp.ClientSession() as session:
        tasks = []

        # Function to run serpapi searches in a separate thread
        def run_search(term):
            print(f"Searching for {term}...")
            search = GoogleSearch({"api_key": API_KEY, "engine": "google", "q": term, "num": MAX_RESULTS})
            results = search.get_dict()
            links = [link['link'] for link in results.get('organic_results', [])]
            return links

        # Running the serpapi search synchronously for each term
        for term in search_terms:
            links = run_search(term)

            # Schedule email fetches for each link asynchronously
            for link in links:
                task = fetch_emails(session, link)
                tasks.append(task)

        # Wait for all fetch tasks to complete
        results = await asyncio.gather(*tasks)

        # Process results
        for exact, found in results:
            for email, domains in exact.items():
                if email not in exact_emails:
                    exact_emails[email] = set()
                exact_emails[email].update(domains)

            for email, domains in found.items():
                if email not in found_emails:
                    found_emails[email] = set()
                found_emails[email].update(domains)

    # Write results to CSV files
    write_emails_to_csv('exact_matches.csv', exact_emails)
    write_emails_to_csv('found_emails.csv', found_emails)

    print("Email extraction finished. Results saved to CSV files.")

if __name__ == "__main__":
    asyncio.run(main_async())